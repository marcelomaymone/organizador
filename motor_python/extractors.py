import re
import zipfile
from abc import ABC, abstractmethod

# =====================================================================
# PRINCIPIOS SOLID: OCP (Open/Closed) e LSP (Liskov Substitution)
# Classes de Extracao herdando de uma classe abstrata comum BaseExtractor
# =====================================================================

class BaseExtractor(ABC):
    """Classe base abstrata para todos os extratores de texto."""

    LIMIT_CHARACTERS = 8000  # Heuristica: 1 token ≈ 4 caracteres, limite de 2000 tokens.

    @abstractmethod
    def extract(self, filepath: str) -> str:
        """Extrai o texto do arquivo e retorna ate o limite de 8.000 caracteres.

        Deve levantar excecoes apropriadas (ex: ValueError, IOError) se o arquivo
        estiver corrompido, inacessivel, criptografado ou for ilegivel.
        """
        pass

    def _truncate_text(self, text: str) -> str:
        """Garante que o texto extraido nao exceda a heuristica de limite."""
        if len(text) > self.LIMIT_CHARACTERS:
            return text[:self.LIMIT_CHARACTERS]
        return text


class TxtExtractor(BaseExtractor):
    """Extrator de texto para arquivos de texto plano e estruturados simples."""

    def extract(self, filepath: str) -> str:
        content = ""
        # Tenta ler como UTF-8
        try:
            with open(filepath, encoding='utf-8', errors='strict') as f:
                content = f.read()
        except UnicodeDecodeError:
            # Fallback para Latin-1 (comum em codificacoes legadas do Windows)
            with open(filepath, encoding='latin-1', errors='ignore') as f:
                content = f.read()
        except Exception as e:
            raise OSError(f"Falha de leitura do arquivo texto: {e}")

        # Limpeza simples de tags HTML/XML
        content = re.sub(r'<[^>]+>', ' ', content)
        # Substitui multiplos espacos por espaco simples
        content = re.sub(r'\s+', ' ', content).strip()

        return self._truncate_text(content)


class PdfExtractor(BaseExtractor):
    """Extrator de texto para arquivos PDF usando PyMuPDF."""

    def extract(self, filepath: str) -> str:
        import fitz  # PyMuPDF

        text_parts = []
        try:
            doc = fitz.open(filepath)
            # Verifica se o PDF esta criptografado
            if doc.is_encrypted:
                raise ValueError("O arquivo PDF esta criptografado.")

            # Processa ate as 3 primeiras paginas
            paginas_limite = min(3, len(doc))
            for i in range(paginas_limite):
                page = doc.load_page(i)
                page_text = page.get_text()
                if page_text:
                    text_parts.append(page_text)
            doc.close()
        except Exception as e:
            raise OSError(f"Falha ao processar o arquivo PDF: {e}")

        full_text = " ".join(text_parts)
        full_text = re.sub(r'\s+', ' ', full_text).strip()

        if not full_text:
            raise ValueError("O PDF nao contem texto extraivel (pode ser composto apenas de imagens).")

        return self._truncate_text(full_text)


class DocxExtractor(BaseExtractor):
    """Extrator de texto para arquivos Word .docx usando python-docx."""

    def extract(self, filepath: str) -> str:
        from docx import Document

        text_parts = []
        try:
            doc = Document(filepath)

            # Extrai texto dos paragrafo
            for paragraph in doc.paragraphs:
                if paragraph.text:
                    text_parts.append(paragraph.text)

            # Extrai texto de tabelas para nao perder metadados tabulares
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text:
                            text_parts.append(cell.text)
        except Exception as e:
            raise OSError(f"Falha ao processar o arquivo DOCX: {e}")

        full_text = " ".join(text_parts)
        full_text = re.sub(r'\s+', ' ', full_text).strip()
        return self._truncate_text(full_text)


class XlsxExtractor(BaseExtractor):
    """Extrator de texto para planilhas Excel .xlsx usando openpyxl."""

    def extract(self, filepath: str) -> str:
        import openpyxl

        text_parts = []
        try:
            # Carrega a planilha em modo de apenas leitura e com formulas calculadas
            wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    for val in row:
                        if val is not None:
                            text_parts.append(str(val))
            wb.close()
        except Exception as e:
            raise OSError(f"Falha ao processar o arquivo XLSX: {e}")

        full_text = " ".join(text_parts)
        full_text = re.sub(r'\s+', ' ', full_text).strip()
        return self._truncate_text(full_text)


class PptxExtractor(BaseExtractor):
    """Extrator de texto para apresentacoes PowerPoint .pptx usando python-pptx."""

    def extract(self, filepath: str) -> str:
        from pptx import Presentation

        text_parts = []
        try:
            prs = Presentation(filepath)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text:
                                text_parts.append(paragraph.text)
        except Exception as e:
            raise OSError(f"Falha ao processar o arquivo PPTX: {e}")

        full_text = " ".join(text_parts)
        full_text = re.sub(r'\s+', ' ', full_text).strip()
        return self._truncate_text(full_text)


class OdtExtractor(BaseExtractor):
    """Extrator de texto para arquivos OpenDocument Text .odt usando odfpy."""

    def extract(self, filepath: str) -> str:
        from odf import opendocument, teletype
        from odf.text import P

        try:
            doc = opendocument.load(filepath)
            paragraphs = doc.getElementsByType(P)
            text_parts = [teletype.extractText(p) for p in paragraphs]
        except Exception as e:
            raise OSError(f"Falha ao processar o arquivo ODT: {e}")

        full_text = " ".join(text_parts)
        full_text = re.sub(r'\s+', ' ', full_text).strip()
        return self._truncate_text(full_text)


class OdsExtractor(BaseExtractor):
    """Extrator de texto para arquivos OpenDocument Spreadsheet .ods usando odfpy."""

    def extract(self, filepath: str) -> str:
        from odf import opendocument, teletype
        from odf.table import Table, TableCell, TableRow

        text_parts = []
        try:
            doc = opendocument.load(filepath)
            sheets = doc.getElementsByType(Table)
            for sheet in sheets:
                rows = sheet.getElementsByType(TableRow)
                for row in rows:
                    cells = row.getElementsByType(TableCell)
                    row_text = []
                    for cell in cells:
                        val = teletype.extractText(cell)
                        if val:
                            row_text.append(val)
                    if row_text:
                        text_parts.append(" ".join(row_text))
        except Exception as e:
            raise OSError(f"Falha ao processar o arquivo ODS: {e}")

        full_text = " ".join(text_parts)
        full_text = re.sub(r'\s+', ' ', full_text).strip()
        return self._truncate_text(full_text)


class OdpExtractor(BaseExtractor):
    """Extrator de texto para arquivos OpenDocument Presentation .odp usando odfpy."""

    def extract(self, filepath: str) -> str:
        from odf import opendocument, teletype
        from odf.draw import Frame
        from odf.text import P

        text_parts = []
        try:
            doc = opendocument.load(filepath)
            frames = doc.getElementsByType(Frame)
            for frame in frames:
                paragraphs = frame.getElementsByType(P)
                for p in paragraphs:
                    val = teletype.extractText(p)
                    if val:
                        text_parts.append(val)
        except Exception as e:
            raise OSError(f"Falha ao processar o arquivo ODP: {e}")

        full_text = " ".join(text_parts)
        full_text = re.sub(r'\s+', ' ', full_text).strip()
        return self._truncate_text(full_text)


class RtfExtractor(BaseExtractor):
    """Extrator de texto para arquivos Rich Text Format .rtf usando striprtf."""

    def extract(self, filepath: str) -> str:
        from striprtf.striprtf import rtf_to_text

        try:
            with open(filepath, encoding='utf-8', errors='ignore') as f:
                rtf_content = f.read()
            text = rtf_to_text(rtf_content)
        except Exception as e:
            raise OSError(f"Falha ao processar o arquivo RTF: {e}")

        text = re.sub(r'\s+', ' ', text).strip()
        return self._truncate_text(text)


class MacExtractor(BaseExtractor):
    """Extrator de texto portavel e leve para arquivos iWork do Mac (.pages, .numbers, .key).

    Garante a portabilidade descompactando o arquivo ZIP do documento iWork
    e extraindo as strings legiveis UTF-8 brutas do binario do banco .iwork correspondente.
    """

    def extract(self, filepath: str) -> str:
        if not zipfile.is_zipfile(filepath):
            raise ValueError("O documento iWork nao e um arquivo compactado (ZIP) valido.")

        text_parts = []
        try:
            with zipfile.ZipFile(filepath, 'r') as zf:
                # 1. Verifica formato legado (iWork '09) contendo index.xml
                if 'index.xml' in zf.namelist():
                    xml_content = zf.read('index.xml').decode('utf-8', errors='ignore')
                    # Limpeza das tags XML
                    text = re.sub(r'<[^>]+>', ' ', xml_content)
                    text_parts.append(text)
                else:
                    # 2. Formato moderno: busca arquivos .iwork (Index/Document.iwork, Index/CalculationEngine.iwork, etc.)
                    alvos_arquivos = [
                        'Index/Document.iwork',             # Pages
                        'Index/CalculationEngine.iwork',    # Numbers
                        'Index/Presentation.iwork',         # Keynote
                        'Document.iwork'                    # Formatos alternativos
                    ]

                    encontrou = False
                    for alvo in alvos_arquivos:
                        if alvo in zf.namelist():
                            encontrou = True
                            binary_data = zf.read(alvo)
                            extracted = self._extract_readable_strings(binary_data)
                            if extracted:
                                text_parts.append(extracted)

                    if not encontrou:
                        # Fallback: le qualquer outro arquivo .iwork ou .xml
                        for name in zf.namelist():
                            if name.endswith('.iwork') or name.endswith('.xml'):
                                binary_data = zf.read(name)
                                extracted = self._extract_readable_strings(binary_data)
                                if extracted:
                                    text_parts.append(extracted)
        except Exception as e:
            raise OSError(f"Falha ao descompactar ou processar o arquivo iWork: {e}")

        full_text = " ".join(text_parts)
        full_text = re.sub(r'\s+', ' ', full_text).strip()

        if not full_text:
            raise ValueError("Nao foi possivel extrair nenhuma string legivel do documento iWork.")

        return self._truncate_text(full_text)

    def _extract_readable_strings(self, data: bytes) -> str:
        """Varre o binario extraindo sequencias imprimiveis de bytes legiveis."""
        # Captura sequencias de 4 ou mais bytes imprimiveis ASCII ou acentuados do Latin-1
        pattern = re.compile(rb'[\x20-\x7E\xC0-\xFF]{4,}')
        matches = pattern.findall(data)
        words = []
        for m in matches:
            # Ignora padroes tipicos de strings de sistema ou XML compactados
            if m.startswith(b'</') or m.startswith(b'<?'):
                continue
            try:
                decoded = m.decode('utf-8').strip()
                if len(decoded) > 3 and not decoded.isspace():
                    words.append(decoded)
            except UnicodeDecodeError:
                try:
                    decoded = m.decode('latin-1').strip()
                    if len(decoded) > 3 and not decoded.isspace():
                        words.append(decoded)
                except Exception:
                    continue
        return " ".join(words)


# =====================================================================
# REGISTRY DE EXTRATORES (Mapeamento de Extensoes)
# =====================================================================

ExtractorRegistry: dict[str, type[BaseExtractor]] = {
    '.txt': TxtExtractor,
    '.md': TxtExtractor,
    '.csv': TxtExtractor,
    '.tsv': TxtExtractor,
    '.html': TxtExtractor,
    '.htm': TxtExtractor,
    '.xml': TxtExtractor,
    '.json': TxtExtractor,

    '.pdf': PdfExtractor,

    '.docx': DocxExtractor,
    '.xlsx': XlsxExtractor,
    '.pptx': PptxExtractor,

    '.odt': OdtExtractor,
    '.ods': OdsExtractor,
    '.odp': OdpExtractor,

    '.rtf': RtfExtractor,

    '.pages': MacExtractor,
    '.numbers': MacExtractor,
    '.key': MacExtractor
}
