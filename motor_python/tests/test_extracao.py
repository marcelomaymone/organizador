import os
import shutil
import sqlite3
import zipfile
import pytest
import fitz
from docx import Document
import openpyxl
from pptx import Presentation
from odf.opendocument import OpenDocumentText, OpenDocumentSpreadsheet, OpenDocumentPresentation
from odf.text import P
from odf.table import Table, TableRow, TableCell
from odf.draw import Frame
from extractors import (
    TxtExtractor, PdfExtractor, DocxExtractor, XlsxExtractor, 
    PptxExtractor, OdtExtractor, OdsExtractor, OdpExtractor, 
    RtfExtractor, MacExtractor, ExtractorRegistry
)
from extrator_worker import ExtractWorker

# Fixtures para criacao de diretorios temporarios e banco de dados de teste

@pytest.fixture
def temp_dirs(tmp_path):
    """Retorna dicionario com pastas temporarias para origem, destino e banco."""
    origem = tmp_path / "origem"
    destino = tmp_path / "destino"
    origem.mkdir()
    destino.mkdir()
    
    db_path = str(tmp_path / "teste_extracao.sqlite")
    
    # Cria o banco de dados temporario com o schema completo
    conn = sqlite3.connect(db_path)
    try:
        conn.execute("""
        CREATE TABLE arquivos_processamento (
            uuid TEXT PRIMARY KEY,
            caminho_origem TEXT UNIQUE,
            nome_original TEXT,
            tamanho_bytes INTEGER,
            hash_xxhash TEXT,
            status TEXT,
            data_criacao_sistema INTEGER,
            data_modificacao_sistema INTEGER,
            justificativa_classificacao TEXT,
            eh_duplicado INTEGER,
            motivo_falha TEXT,
            mensagem_erro TEXT,
            texto_extraido TEXT,
            data_registro TEXT
        );
        """)
        conn.commit()
    finally:
        conn.close()
        
    return {
        'origem': origem,
        'destino': destino,
        'db_path': db_path
    }

# =====================================================================
# TESTES UNITARIOS DE EXTRATORES INDIVIDUAIS
# =====================================================================

def test_txt_extractor(temp_dirs):
    """Valida a extracao de texto simples com acentos e remocao de tags HTML."""
    filepath = os.path.join(temp_dirs['origem'], "texto.txt")
    with open(filepath, "w", encoding="utf-8") as f:
        f.write("<h1>Olá Mundo!</h1> Esse é um teste de acentuação em português.")
        
    extractor = TxtExtractor()
    texto = extractor.extract(filepath)
    
    assert "Olá Mundo!" in texto
    assert "acentuação" in texto
    assert "<h1>" not in texto


def test_txt_extractor_truncation(temp_dirs):
    """Valida o truncamento do texto no limite heuristico de 8.000 caracteres."""
    filepath = os.path.join(temp_dirs['origem'], "grande.txt")
    conteudo = "A" * 10000
    with open(filepath, "w", encoding="utf-8") as f:
        f.write(conteudo)
        
    extractor = TxtExtractor()
    texto = extractor.extract(filepath)
    
    assert len(texto) == 8000


def test_pdf_extractor(temp_dirs):
    """Valida a extracao de PDF gerando um documento PDF valido dinamicamente."""
    filepath = os.path.join(temp_dirs['origem'], "teste.pdf")
    
    # Cria um PDF valido usando PyMuPDF
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((50, 50), "Conteudo de teste do PDF em Portugues.")
    doc.save(filepath)
    doc.close()
    
    extractor = PdfExtractor()
    texto = extractor.extract(filepath)
    
    assert "Conteudo de teste" in texto


def test_docx_extractor(temp_dirs):
    """Valida a extracao de DOCX gerando um documento Word valido."""
    filepath = os.path.join(temp_dirs['origem'], "teste.docx")
    
    doc = Document()
    doc.add_paragraph("Parágrafo de teste do Word.")
    # Adiciona tabela
    table = doc.add_table(rows=1, cols=1)
    table.cell(0, 0).text = "Texto na tabela do Word"
    doc.save(filepath)
    
    extractor = DocxExtractor()
    texto = extractor.extract(filepath)
    
    assert "Parágrafo de teste" in texto
    assert "Texto na tabela" in texto


def test_xlsx_extractor(temp_dirs):
    """Valida a extracao de XLSX gerando uma planilha Excel valida."""
    filepath = os.path.join(temp_dirs['origem'], "teste.xlsx")
    
    wb = openpyxl.Workbook()
    ws = wb.active
    ws['A1'] = "Item de teste Excel"
    ws['B2'] = 12345
    wb.save(filepath)
    wb.close()
    
    extractor = XlsxExtractor()
    texto = extractor.extract(filepath)
    
    assert "Item de teste" in texto
    assert "12345" in texto


def test_pptx_extractor(temp_dirs):
    """Valida a extracao de PPTX gerando uma apresentacao PowerPoint valida."""
    filepath = os.path.join(temp_dirs['origem'], "teste.pptx")
    
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Layout em branco
    tx_box = slide.shapes.add_textbox(100, 100, 200, 100)
    tf = tx_box.text_frame
    tf.text = "Texto de teste do slide PowerPoint"
    prs.save(filepath)
    
    extractor = PptxExtractor()
    texto = extractor.extract(filepath)
    
    assert "Texto de teste" in texto


def test_rtf_extractor(temp_dirs):
    """Valida a extracao de RTF."""
    filepath = os.path.join(temp_dirs['origem'], "teste.rtf")
    with open(filepath, "w", encoding="ascii") as f:
        f.write(r"{\rtf1\ansi\deff0 {\fonttbl {\f0 Courier;}}\f0\fs24 Texto limpo RTF.}")
        
    extractor = RtfExtractor()
    texto = extractor.extract(filepath)
    
    assert "Texto limpo RTF" in texto


def test_odf_extractors(temp_dirs):
    """Valida a extracao de ODT, ODS e ODP."""
    odt_path = os.path.join(temp_dirs['origem'], "teste.odt")
    ods_path = os.path.join(temp_dirs['origem'], "teste.ods")
    odp_path = os.path.join(temp_dirs['origem'], "teste.odp")
    
    # 1. ODT
    odt = OpenDocumentText()
    odt.text.addElement(P(text="Conteudo ODT OpenOffice"))
    odt.save(odt_path)
    
    # 2. ODS
    ods = OpenDocumentSpreadsheet()
    table = Table(name="Planilha1")
    row = TableRow()
    cell = TableCell(valuetype="string")
    cell.addElement(P(text="Celula ODS"))
    row.addElement(cell)
    table.addElement(row)
    ods.spreadsheet.addElement(table)
    ods.save(ods_path)
    
    # 3. ODP
    odp = OpenDocumentPresentation()
    # Adiciona elemento em Page e depois na presentation
    from odf.draw import Page, TextBox
    page = Page(name="Slide1", masterpagename="Master1")
    frame = Frame()
    tb = TextBox()
    tb.addElement(P(text="Texto ODP Apresentacao"))
    frame.addElement(tb)
    page.addElement(frame)
    odp.presentation.addElement(page)
    odp.save(odp_path)
    
    assert "ODT" in OdtExtractor().extract(odt_path)
    assert "Celula ODS" in OdsExtractor().extract(ods_path)
    assert "Texto ODP" in OdpExtractor().extract(odp_path)


def test_mac_extractor(temp_dirs):
    """Valida o MacExtractor simulando um arquivo zip Pages do iWork contendo binario .iwork."""
    filepath = os.path.join(temp_dirs['origem'], "doc.pages")
    
    # Cria zip simulado contendo o binario Index/Document.iwork
    with zipfile.ZipFile(filepath, 'w') as zf:
        zf.writestr('Index/Document.iwork', b'\x00\x01\x03\x02Texto de Teste do Mac Pages\x00\x00')
        
    extractor = MacExtractor()
    texto = extractor.extract(filepath)
    
    assert "Texto de Teste do Mac Pages" in texto

# =====================================================================
# TESTES INTEGRADOS DO WORKER (FILA, QUARENTENA E PROCESSAMENTO)
# =====================================================================

def test_extract_worker_success(temp_dirs):
    """Garante que o worker processa arquivos validos na fila com sucesso."""
    # Cria arquivo de teste
    f1 = os.path.join(temp_dirs['origem'], "doc1.txt")
    with open(f1, "w", encoding="utf-8") as f:
        f.write("Conteudo semantico de teste")
        
    # Insere na fila
    conn = sqlite3.connect(temp_dirs['db_path'])
    try:
        conn.execute(
            """
            INSERT INTO arquivos_processamento 
            (uuid, caminho_origem, nome_original, tamanho_bytes, hash_xxhash, status, eh_duplicado)
            VALUES ('uuid-ok', ?, 'doc1.txt', 10, 'hash1', 'pendente_extracao', 0)
            """,
            (f1,)
        )
        conn.commit()
    finally:
        conn.close()
        
    worker = ExtractWorker(temp_dirs['db_path'], str(temp_dirs['destino']))
    total = worker.execute()
    
    assert total == 1
    
    # Valida no banco
    conn = sqlite3.connect(temp_dirs['db_path'])
    try:
        cursor = conn.cursor()
        cursor.execute("SELECT status, texto_extraido, motivo_falha FROM arquivos_processamento WHERE uuid='uuid-ok'")
        row = cursor.fetchone()
        assert row is not None
        assert row[0] == 'pendente_inferencia'
        assert row[1] == "Conteudo semantico de teste"
        assert row[2] is None
    finally:
        conn.close()


def test_extract_worker_unsupported_format(temp_dirs):
    """Garante que arquivos nao-suportados avancam de status com justificativa automatica."""
    f1 = os.path.join(temp_dirs['origem'], "programa.exe")
    with open(f1, "wb") as f:
        f.write(b"\x90\x00\x00\x00")
        
    conn = sqlite3.connect(temp_dirs['db_path'])
    try:
        conn.execute(
            """
            INSERT INTO arquivos_processamento 
            (uuid, caminho_origem, nome_original, tamanho_bytes, hash_xxhash, status, eh_duplicado)
            VALUES ('uuid-exe', ?, 'programa.exe', 4, 'hash2', 'pendente_extracao', 0)
            """,
            (f1,)
        )
        conn.commit()
    finally:
        conn.close()
        
    worker = ExtractWorker(temp_dirs['db_path'], str(temp_dirs['destino']))
    total = worker.execute()
    
    assert total == 1
    
    conn = sqlite3.connect(temp_dirs['db_path'])
    try:
        cursor = conn.cursor()
        cursor.execute("SELECT status, texto_extraido, justificativa_classificacao FROM arquivos_processamento WHERE uuid='uuid-exe'")
        row = cursor.fetchone()
        assert row is not None
        assert row[0] == 'pendente_inferencia'
        assert row[1] == ''
        assert "Formato de arquivo não textual" in row[2]
    finally:
        conn.close()


def test_extract_worker_quarantine_flow(temp_dirs):
    """Garante que arquivos corrompidos sao isolados fisicamente na quarentena absoluta (Opcao B)."""
    # Cria PDF corrompido
    f1 = os.path.join(temp_dirs['origem'], "corrompido.pdf")
    with open(f1, "w", encoding="utf-8") as f:
        f.write("dados invalidos pdf")
        
    conn = sqlite3.connect(temp_dirs['db_path'])
    try:
        conn.execute(
            """
            INSERT INTO arquivos_processamento 
            (uuid, caminho_origem, nome_original, tamanho_bytes, hash_xxhash, status, eh_duplicado)
            VALUES ('uuid-corrompido', ?, 'corrompido.pdf', 20, 'hash3', 'pendente_extracao', 0)
            """,
            (f1,)
        )
        conn.commit()
    finally:
        conn.close()
        
    worker = ExtractWorker(temp_dirs['db_path'], str(temp_dirs['destino']))
    total = worker.execute()
    
    assert total == 1
    
    # 1. Valida se o arquivo original foi movido/removido da origem
    assert not os.path.exists(f1)
    
    # 2. Valida se a quarentena fisica foi criada no destino preservando o caminho do drive absoluto (Opcao B)
    drive, resto = os.path.splitdrive(f1)
    drive_limpo = drive.replace(":", "_").strip("\\/")
    resto_limpo = resto.lstrip("\\/")
    
    caminho_quarentena_esperado = os.path.join(
        temp_dirs['destino'], 
        "_QUARENTENA_", 
        drive_limpo, 
        resto_limpo
    )
    
    caminho_verificar_exists = caminho_quarentena_esperado
    if os.name == 'nt' and len(caminho_quarentena_esperado) > 240:
        caminho_verificar_exists = "\\\\?\\" + os.path.abspath(caminho_quarentena_esperado)
        
    assert os.path.exists(caminho_verificar_exists)
    
    # 3. Valida no banco se o status mudou para 'quarentena', o caminho_origem atualizou para a quarentena e gravou o motivo
    conn = sqlite3.connect(temp_dirs['db_path'])
    try:
        cursor = conn.cursor()
        cursor.execute("SELECT status, caminho_origem, motivo_falha, mensagem_erro FROM arquivos_processamento WHERE uuid='uuid-corrompido'")
        row = cursor.fetchone()
        assert row is not None
        assert row[0] == 'quarentena'
        assert row[1] == caminho_quarentena_esperado
        assert "Falha na extração de texto" in row[2]
        assert row[3] is not None
    finally:
        conn.close()
